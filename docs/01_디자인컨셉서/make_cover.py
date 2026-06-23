# -*- coding: utf-8 -*-
"""디자인컨셉서 표지(히어로) 슬라이드를 PPTX로 생성."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import copy

EMU = 914400

# ---- helpers ---------------------------------------------------------------

def set_grad_bg(shape, stops):
    """shape 채우기를 선형 그라데이션으로. stops=[(pos0~100, 'RRGGBB'), ...], angle 145deg."""
    sp = shape.fill._xPr  # spPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for e in sp.findall(qn(tag)):
            sp.remove(e)
    grad = etree.SubElement(sp, qn('a:gradFill'))
    lst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, hexc in stops:
        gs = etree.SubElement(lst, qn('a:gs'))
        gs.set('pos', str(int(pos * 1000)))
        c = etree.SubElement(gs, qn('a:srgbClr'))
        c.set('val', hexc)
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(int(145 * 60000)))  # 145deg
    lin.set('scaled', '1')
    # ensure correct child order: spPr expects gradFill before a:ln if present
    ln = sp.find(qn('a:ln'))
    if ln is not None:
        sp.remove(grad)
        ln.addprevious(grad)


def set_solid_alpha(shape, hexc, alpha_pct):
    """단색 + 투명도(알파, 0~100 불투명도)."""
    sp = shape.fill._xPr
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for e in sp.findall(qn(tag)):
            sp.remove(e)
    solid = etree.SubElement(sp, qn('a:solidFill'))
    c = etree.SubElement(solid, qn('a:srgbClr'))
    c.set('val', hexc)
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))
    ln = sp.find(qn('a:ln'))
    if ln is not None:
        sp.remove(solid)
        ln.addprevious(solid)


def no_line(shape):
    shape.line.fill.background()


def grad_text(run, stops):
    """run 글자에 선형 그라데이션 채우기."""
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill'):
        for e in rPr.findall(qn(tag)):
            rPr.remove(e)
    grad = etree.Element(qn('a:gradFill'))
    lst = etree.SubElement(grad, qn('a:gsLst'))
    for pos, hexc in stops:
        gs = etree.SubElement(lst, qn('a:gs'))
        gs.set('pos', str(int(pos * 1000)))
        c = etree.SubElement(gs, qn('a:srgbClr'))
        c.set('val', hexc)
    lin = etree.SubElement(grad, qn('a:lin'))
    lin.set('ang', str(int(125 * 60000)))
    lin.set('scaled', '1')
    # rPr children order: fills go after a:ln, before a:latin etc. Insert at correct spot.
    ln = rPr.find(qn('a:ln'))
    if ln is not None:
        ln.addnext(grad)
    else:
        rPr.insert(0, grad)


def style_run(run, font='Pretendard', size=None, bold=True, color=None,
              spacing=None, letter_spacing=None):
    f = run.font
    f.name = font
    if size is not None:
        f.size = Pt(size)
    f.bold = bold
    if color is not None:
        f.color.rgb = RGBColor.from_string(color)
    rPr = run._r.get_or_add_rPr()
    # east asian + complex script 폰트도 Pretendard로
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(rPr, qn(tag))
        el.set('typeface', font)
    if letter_spacing is not None:
        rPr.set('spc', str(int(letter_spacing)))


# ---- build -----------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# 1) 배경 이미지 (다크 네이비 그라데이션 + 글로우, 풀블리드)
slide.shapes.add_picture(
    '/Users/gosujin/klid_2nd/docs/01_디자인컨셉서/표지_배경.png',
    0, 0, SW, SH)

LEFT = Inches(0.98)
WIDTH = Inches(11.4)

# 3) eyebrow 라인 + 텍스트
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT, Inches(1.62),
                              Inches(0.30), Pt(2.2))
no_line(line)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor.from_string('AEBCFF')
line.shadow.inherit = False

eb = slide.shapes.add_textbox(LEFT + Inches(0.42), Inches(1.40),
                              WIDTH, Inches(0.5))
tf = eb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = 'HERO DESIGN CONCEPT'
style_run(r, size=13, bold=True, color='AEBCFF', letter_spacing=240)

# 4) 타이틀 H1
title = slide.shapes.add_textbox(LEFT, Inches(2.05), WIDTH, Inches(3.0))
tf = title.text_frame
tf.word_wrap = True
# 1줄: KLID
p1 = tf.paragraphs[0]
p1.line_spacing = 0.96
r = p1.add_run(); r.text = 'KLID'
style_run(r, size=80, bold=True, color='FFFFFF', letter_spacing=-300)
# 2줄: AI 영상학습 사용자 포털 (그라데이션)
p2 = tf.add_paragraph()
p2.line_spacing = 0.98
r = p2.add_run(); r.text = 'AI 영상학습 사용자 포털'
style_run(r, size=66, bold=True, letter_spacing=-300)
grad_text(r, [(0, 'FFFFFF'), (52, 'B9C6FF'), (100, '8FF4FF')])

# 5) 서브 카피
sub = slide.shapes.add_textbox(LEFT, Inches(5.55), Inches(10.2), Inches(1.4))
tf = sub.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.line_spacing = 1.5
r = p.add_run()
r.text = '데이터를 찾고, 키우고, 바로 쓸 수 있게 만드는 AI 영상학습 사용자 포털의 브랜드 경험 방향.'
style_run(r, size=18, bold=False, color='C2CCE6')

# 6) 하단 디바이더
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, LEFT, Inches(6.92),
                             Inches(11.37), Pt(1))
no_line(div)
set_solid_alpha(div, 'FFFFFF', 22)
div.shadow.inherit = False

out = '/Users/gosujin/klid_2nd/docs/01_디자인컨셉서/디자인컨셉서_표지.pptx'
prs.save(out)
print('saved:', out)
